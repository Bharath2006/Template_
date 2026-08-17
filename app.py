import os
import json
import uuid
import shutil
import tempfile
from datetime import datetime
from flask import Flask, request, jsonify, send_file, session
from flask_cors import CORS
from werkzeug.utils import secure_filename
from dotenv import load_dotenv
import google.generativeai as genai

# Document processing libraries
from pptx import Presentation
from docx import Document
from openpyxl import load_workbook
import fitz  # PyMuPDF
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from PIL import Image
import io

load_dotenv()

app = Flask(__name__)
app.secret_key = os.getenv('SECRET_KEY', 'dev-secret-key')
CORS(app)

# Configuration
UPLOAD_FOLDER = tempfile.mkdtemp()
OUTPUT_FOLDER = tempfile.mkdtemp()
ALLOWED_EXTENSIONS = {'pptx', 'pdf', 'docx', 'xlsx'}
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB

# Gemini configuration
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
if not GEMINI_API_KEY:
    raise ValueError("GEMINI_API_KEY not found in environment variables")

genai.configure(api_key=GEMINI_API_KEY)
gemini_model = genai.GenerativeModel('gemini-1.5-flash')

# Store uploaded files and processed data
file_store = {}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_pptx_structure(filepath):
    """Extract structure from PPTX file"""
    prs = Presentation(filepath)
    structure = {
        'type': 'pptx',
        'slides': []
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_data = {
            'slide_number': slide_idx + 1,
            'elements': []
        }
        
        for shape in slide.shapes:
            element = {
                'type': str(shape.shape_type),
                'name': shape.name,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            }
            
            if hasattr(shape, 'text'):
                element['text'] = shape.text
                element['text_frame'] = {
                    'text': shape.text,
                    'paragraphs': []
                }
                for paragraph in shape.text_frame.paragraphs:
                    para_data = {
                        'text': paragraph.text,
                        'alignment': str(paragraph.alignment) if paragraph.alignment else None,
                        'font_size': paragraph.font.size if paragraph.font.size else None,
                        'font_name': paragraph.font.name,
                        'bold': paragraph.font.bold,
                        'italic': paragraph.font.italic,
                        'underline': paragraph.font.underline
                    }
                    element['text_frame']['paragraphs'].append(para_data)
            
            if hasattr(shape, 'fill') and hasattr(shape.fill, 'fore_color'):
                try:
                    element['fill_color'] = str(shape.fill.fore_color.rgb)
                except:
                    pass
            
            slide_data['elements'].append(element)
        
        structure['slides'].append(slide_data)
    
    return structure

def extract_pdf_structure(filepath):
    """Extract structure from PDF file"""
    doc = fitz.open(filepath)
    structure = {
        'type': 'pdf',
        'pages': []
    }
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        page_data = {
            'page_number': page_num + 1,
            'width': page.rect.width,
            'height': page.rect.height,
            'text_blocks': []
        }
        
        text_blocks = page.get_text("dict")
        for block in text_blocks.get('blocks', []):
            if 'lines' in block:
                for line in block['lines']:
                    for span in line.get('spans', []):
                        text_block = {
                            'text': span.get('text', ''),
                            'font': span.get('font', ''),
                            'size': span.get('size', 0),
                            'color': span.get('color', 0),
                            'origin': span.get('origin', [0, 0])
                        }
                        if text_block['text'].strip():
                            page_data['text_blocks'].append(text_block)
        
        # Extract images
        images = page.get_images()
        page_data['image_count'] = len(images)
        
        structure['pages'].append(page_data)
    
    doc.close()
    return structure

def extract_docx_structure(filepath):
    """Extract structure from DOCX file"""
    doc = Document(filepath)
    structure = {
        'type': 'docx',
        'paragraphs': [],
        'tables': []
    }
    
    for para in doc.paragraphs:
        para_data = {
            'text': para.text,
            'style': para.style.name if para.style else None,
            'alignment': str(para.alignment) if para.alignment else None
        }
        
        # Extract runs for formatting
        runs_data = []
        for run in para.runs:
            run_info = {
                'text': run.text,
                'bold': run.bold,
                'italic': run.italic,
                'underline': run.underline,
                'font_name': run.font.name if run.font and run.font.name else None,
                'font_size': run.font.size if run.font and run.font.size else None
            }
            runs_data.append(run_info)
        para_data['runs'] = runs_data
        structure['paragraphs'].append(para_data)
    
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_text = cell.text
                row_data.append(cell_text)
            table_data.append(row_data)
        structure['tables'].append(table_data)
    
    return structure

def extract_xlsx_structure(filepath):
    """Extract structure from XLSX file"""
    wb = load_workbook(filepath, data_only=True)
    structure = {
        'type': 'xlsx',
        'sheets': []
    }
    
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_data = {
            'name': sheet_name,
            'rows': [],
            'columns': []
        }
        
        # Get data
        for row in ws.iter_rows(values_only=True):
            row_data = [cell if cell is not None else '' for cell in row]
            if any(row_data):
                sheet_data['rows'].append(row_data)
        
        # Get column widths
        for col in ws.columns:
            col_letter = col[0].column_letter
            col_width = ws.column_dimensions[col_letter].width
            sheet_data['columns'].append({
                'letter': col_letter,
                'width': col_width if col_width else 10
            })
        
        structure['sheets'].append(sheet_data)
    
    return structure

def apply_pptx_template(template_structure, content_data, output_path):
    """Apply content to PPTX template"""
    prs = Presentation()
    
    # Use template structure to create presentation
    for slide_idx, slide_template in enumerate(template_structure.get('slides', [])):
        slide_layout = prs.slide_layouts[0]  # Use first layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Get content for this slide from Gemini response
        slide_content = content_data.get('slides', [{}])[slide_idx] if slide_idx < len(content_data.get('slides', [])) else {}
        
        # Add content to slide
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                # Find matching content
                content_key = shape.name or f"slide_{slide_idx}_text"
                new_text = slide_content.get(content_key, shape.text if hasattr(shape, 'text') else '')
                if new_text:
                    shape.text = str(new_text)
    
    prs.save(output_path)
    return output_path

def apply_pdf_template(template_structure, content_data, output_path):
    """Apply content to PDF template"""
    doc = SimpleDocTemplate(output_path, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    # Add content from template structure and content data
    pages = template_structure.get('pages', [])
    content_pages = content_data.get('pages', [])
    
    for page_idx, page in enumerate(pages):
        # Add heading
        story.append(Paragraph(f"Page {page_idx + 1}", styles['Heading1']))
        story.append(Spacer(1, 0.2 * inch))
        
        # Add text blocks
        text_blocks = page.get('text_blocks', [])
        content_text_blocks = content_pages[page_idx].get('text_blocks', []) if page_idx < len(content_pages) else []
        
        for block_idx, block in enumerate(text_blocks):
            text = block.get('text', '')
            if text.strip():
                # Use content if available
                if block_idx < len(content_text_blocks):
                    text = content_text_blocks[block_idx].get('text', text)
                story.append(Paragraph(text, styles['Normal']))
                story.append(Spacer(1, 0.1 * inch))
    
    doc.build(story)
    return output_path

def apply_docx_template(template_structure, content_data, output_path):
    """Apply content to DOCX template"""
    doc = Document()
    
    # Add content based on template structure
    paragraphs = template_structure.get('paragraphs', [])
    content_paragraphs = content_data.get('paragraphs', [])
    
    for idx, para_template in enumerate(paragraphs):
        if idx < len(content_paragraphs):
            text = content_paragraphs[idx].get('text', para_template.get('text', ''))
        else:
            text = para_template.get('text', '')
        
        if text.strip():
            doc.add_paragraph(text)
    
    # Add tables
    tables = template_structure.get('tables', [])
    content_tables = content_data.get('tables', [])
    
    for table_idx, table_template in enumerate(tables):
        if table_idx < len(content_tables):
            table_data = content_tables[table_idx]
            table = doc.add_table(rows=len(table_data), cols=len(table_data[0]) if table_data else 1)
            
            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_data in enumerate(row_data):
                    if col_idx < len(table.rows[row_idx].cells):
                        table.rows[row_idx].cells[col_idx].text = str(cell_data)
    
    doc.save(output_path)
    return output_path

def apply_xlsx_template(template_structure, content_data, output_path):
    """Apply content to XLSX template"""
    from openpyxl import Workbook
    
    wb = Workbook()
    
    sheets = template_structure.get('sheets', [])
    content_sheets = content_data.get('sheets', [])
    
    for sheet_idx, sheet_template in enumerate(sheets):
        if sheet_idx == 0:
            ws = wb.active
            ws.title = sheet_template.get('name', 'Sheet1')
        else:
            ws = wb.create_sheet(sheet_template.get('name', f'Sheet{sheet_idx + 1}'))
        
        # Add data
        rows = sheet_template.get('rows', [])
        content_rows = content_sheets[sheet_idx].get('rows', []) if sheet_idx < len(content_sheets) else []
        
        for row_idx, row_data in enumerate(rows):
            if row_idx < len(content_rows):
                data = content_rows[row_idx]
            else:
                data = row_data
            
            for col_idx, value in enumerate(data):
                ws.cell(row=row_idx + 1, column=col_idx + 1, value=value)
        
        # Set column widths
        for col_idx, col_data in enumerate(sheet_template.get('columns', [])):
            col_letter = col_data.get('letter')
            if col_letter:
                ws.column_dimensions[col_letter].width = col_data.get('width', 10)
    
    wb.save(output_path)
    return output_path

def render_output(file_type, template_structure, content_data):
    """Render output based on file type"""
    file_id = str(uuid.uuid4())
    output_path = os.path.join(OUTPUT_FOLDER, f"{file_id}.{file_type}")
    
    if file_type == 'pptx':
        apply_pptx_template(template_structure, content_data, output_path)
    elif file_type == 'pdf':
        apply_pdf_template(template_structure, content_data, output_path)
    elif file_type == 'docx':
        apply_docx_template(template_structure, content_data, output_path)
    elif file_type == 'xlsx':
        apply_xlsx_template(template_structure, content_data, output_path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
    
    return file_id, output_path

@app.route('/api/health', methods=['GET'])
def health_check():
    return jsonify({'status': 'healthy', 'timestamp': datetime.now().isoformat()})

@app.route('/api/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': 'File type not supported'}), 400
    
    # Check file size
    file.seek(0, os.SEEK_END)
    file_size = file.tell()
    file.seek(0)
    if file_size > MAX_FILE_SIZE:
        return jsonify({'error': 'File size exceeds 50MB limit'}), 400
    
    # Save file
    file_id = str(uuid.uuid4())
    filename = secure_filename(file.filename)
    file_ext = filename.rsplit('.', 1)[1].lower()
    saved_path = os.path.join(UPLOAD_FOLDER, f"{file_id}.{file_ext}")
    file.save(saved_path)
    
    file_store[file_id] = {
        'id': file_id,
        'filename': filename,
        'path': saved_path,
        'type': file_ext,
        'uploaded_at': datetime.now().isoformat(),
        'structure': None,
        'content': None
    }
    
    return jsonify({
        'file_id': file_id,
        'filename': filename,
        'type': file_ext,
        'message': 'File uploaded successfully'
    })

@app.route('/api/analyze', methods=['POST'])
def analyze_file():
    data = request.json
    file_id = data.get('file_id')
    
    if not file_id or file_id not in file_store:
        return jsonify({'error': 'File not found'}), 404
    
    file_info = file_store[file_id]
    file_path = file_info['path']
    file_type = file_info['type']
    
    # Extract structure based on file type
    try:
        if file_type == 'pptx':
            structure = extract_pptx_structure(file_path)
        elif file_type == 'pdf':
            structure = extract_pdf_structure(file_path)
        elif file_type == 'docx':
            structure = extract_docx_structure(file_path)
        elif file_type == 'xlsx':
            structure = extract_xlsx_structure(file_path)
        else:
            return jsonify({'error': 'Unsupported file type'}), 400
        
        file_info['structure'] = structure
        return jsonify({
            'file_id': file_id,
            'type': file_type,
            'structure': structure,
            'message': 'Analysis complete'
        })
    except Exception as e:
        return jsonify({'error': f'Analysis failed: {str(e)}'}), 500

@app.route('/api/generate', methods=['POST'])
def generate_content():
    data = request.json
    file_id = data.get('file_id')
    user_prompt = data.get('prompt', '')
    
    if not file_id or file_id not in file_store:
        return jsonify({'error': 'File not found'}), 404
    
    file_info = file_store[file_id]
    structure = file_info.get('structure')
    
    if not structure:
        return jsonify({'error': 'Please analyze the file first'}), 400
    
    # Prepare prompt for Gemini
    gemini_prompt = f"""
    You are a document/content generation expert. Given a template structure and user requirements, 
    generate appropriate content that fits perfectly into the existing template structure.
    
    TEMPLATE STRUCTURE (JSON):
    {json.dumps(structure, indent=2)}
    
    USER REQUIREMENTS:
    {user_prompt}
    
    INSTRUCTIONS:
    1. Generate content that fits exactly into the existing structure
    2. Preserve the original layout, positioning, and styling
    3. Return ONLY valid JSON that matches the structure
    4. The JSON should have the same keys and structure as the template
    5. Replace placeholder text with meaningful content based on the user requirements
    6. Do not add new elements that don't exist in the template
    7. For text elements, generate appropriate text content
    8. Return ONLY the JSON, no additional text or code
    9. Ensure the JSON is valid and parseable
    
    Response (JSON only):
    """
    
    try:
        # Generate content using Gemini
        response = gemini_model.generate_content(gemini_prompt)
        content_json = response.text
        
        # Clean the response to get valid JSON
        content_json = content_json.strip()
        if content_json.startswith('```json'):
            content_json = content_json[7:]
        if content_json.startswith('```'):
            content_json = content_json[3:]
        if content_json.endswith('```'):
            content_json = content_json[:-3]
        content_json = content_json.strip()
        
        # Parse JSON
        parsed_content = json.loads(content_json)
        
        file_info['content'] = parsed_content
        
        return jsonify({
            'file_id': file_id,
            'content': parsed_content,
            'message': 'Content generated successfully'
        })
    except json.JSONDecodeError as e:
        return jsonify({'error': f'Invalid JSON from Gemini: {str(e)}'}), 500
    except Exception as e:
        return jsonify({'error': f'Generation failed: {str(e)}'}), 500

@app.route('/api/validate', methods=['POST'])
def validate_content():
    data = request.json
    file_id = data.get('file_id')
    content = data.get('content')
    
    if not file_id or file_id not in file_store:
        return jsonify({'error': 'File not found'}), 404
    
    file_info = file_store[file_id]
    structure = file_info.get('structure')
    
    if not structure:
        return jsonify({'error': 'Template structure not found'}), 400
    
    # Validate that content matches structure
    try:
        # Simple structural validation
        if not isinstance(content, dict):
            return jsonify({'error': 'Content must be a JSON object'}), 400
        
        # Validate required keys
        if 'type' in structure and structure['type'] != content.get('type'):
            return jsonify({'error': f"Content type mismatch. Expected {structure['type']}"}), 400
        
        # Validate slides/pages/paragraphs count
        if 'slides' in structure and 'slides' in content:
            if len(structure['slides']) != len(content['slides']):
                return jsonify({'error': f"Slide count mismatch. Expected {len(structure['slides'])}, got {len(content['slides'])}"}), 400
        
        file_info['content'] = content
        
        return jsonify({
            'file_id': file_id,
            'valid': True,
            'message': 'Validation successful'
        })
    except Exception as e:
        return jsonify({'error': f'Validation failed: {str(e)}'}), 500

@app.route('/api/apply', methods=['POST'])
def apply_content():
    data = request.json
    file_id = data.get('file_id')
    
    if not file_id or file_id not in file_store:
        return jsonify({'error': 'File not found'}), 404
    
    file_info = file_store[file_id]
    structure = file_info.get('structure')
    content = file_info.get('content')
    
    if not structure or not content:
        return jsonify({'error': 'Missing structure or content'}), 400
    
    file_info['applied'] = True
    
    return jsonify({
        'file_id': file_id,
        'message': 'Content applied successfully'
    })

@app.route('/api/render', methods=['POST'])
def render_document():
    data = request.json
    file_id = data.get('file_id')
    output_format = data.get('format')
    
    if not file_id or file_id not in file_store:
        return jsonify({'error': 'File not found'}), 404
    
    if not output_format or output_format not in ['pptx', 'pdf', 'docx', 'xlsx']:
        return jsonify({'error': 'Invalid output format'}), 400
    
    file_info = file_store[file_id]
    structure = file_info.get('structure')
    content = file_info.get('content')
    
    if not structure or not content:
        return jsonify({'error': 'Missing structure or content'}), 400
    
    try:
        # Render the output
        output_file_id, output_path = render_output(
            output_format,
            structure,
            content
        )
        
        file_info['output_file_id'] = output_file_id
        file_info['output_path'] = output_path
        
        return jsonify({
            'file_id': file_id,
            'output_file_id': output_file_id,
            'format': output_format,
            'message': 'Rendering complete'
        })
    except Exception as e:
        return jsonify({'error': f'Rendering failed: {str(e)}'}), 500

@app.route('/api/download/<file_id>', methods=['GET'])
def download_file(file_id):
    # Check if this is an output file
    for file_info in file_store.values():
        if file_info.get('output_file_id') == file_id:
            output_path = file_info.get('output_path')
            if output_path and os.path.exists(output_path):
                return send_file(
                    output_path,
                    as_attachment=True,
                    download_name=f"generated_output.{file_info.get('type', 'pdf')}"
                )
    
    return jsonify({'error': 'File not found'}), 404

@app.route('/api/file/<file_id>', methods=['DELETE'])
def delete_file(file_id):
    if file_id in file_store:
        file_info = file_store[file_id]
        if os.path.exists(file_info.get('path', '')):
            os.remove(file_info['path'])
        if os.path.exists(file_info.get('output_path', '')):
            os.remove(file_info['output_path'])
        del file_store[file_id]
        return jsonify({'message': 'File deleted successfully'})
    return jsonify({'error': 'File not found'}), 404

if __name__ == '__main__':
    app.run(debug=True, port=5000)
